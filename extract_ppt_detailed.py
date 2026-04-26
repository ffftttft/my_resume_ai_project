"""Extract detailed content from reference PPT to understand logic flow."""
import sys
import io
from pptx import Presentation
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

ref_path = Path(r"C:\Users\18153\Desktop\人工智能应用_马莹奇_总文件夹\01_作品与答辩材料\慧眼识途 - 基于无监督学习的智能交通异常预警系统.pptx")

prs = Presentation(ref_path)
print(f"总页数: {len(prs.slides)}\n")
print("=" * 80)

content_pages = []

for i, slide in enumerate(prs.slides, 1):
    title = slide.shapes.title.text if slide.shapes.title else "(无标题)"

    # Extract all text content
    all_texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            if text and text != title:
                all_texts.append(text)

    # Only show slides with substantial content
    if all_texts or title != "(无标题)":
        print(f"\n【第 {i} 页】 {title}")
        print("-" * 80)
        for text in all_texts:
            # Show full text for content pages
            if len(text) < 300:
                print(f"  {text}")
            else:
                print(f"  {text[:300]}...")

        if all_texts:
            content_pages.append(i)

print("\n" + "=" * 80)
print(f"\n内容页编号: {content_pages}")
print(f"实际内容页数: {len(content_pages)}")
