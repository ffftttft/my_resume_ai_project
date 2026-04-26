"""Extract structure from reference PPT."""
import sys
import io
from pptx import Presentation
from pathlib import Path

# Fix Windows console encoding
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

ref_path = Path(r"C:\Users\18153\Desktop\人工智能应用_马莹奇_总文件夹\01_作品与答辩材料\慧眼识途 - 基于无监督学习的智能交通异常预警系统.pptx")

prs = Presentation(ref_path)
print(f"总页数: {len(prs.slides)}\n")
print("=" * 60)

for i, slide in enumerate(prs.slides, 1):
    title = slide.shapes.title.text if slide.shapes.title else "(无标题)"
    print(f"\n第 {i} 页: {title}")
    print("-" * 60)

    # Extract text from all shapes
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            if shape != slide.shapes.title:
                texts.append(shape.text.strip())

    if texts:
        for text in texts[:3]:  # Show first 3 text blocks
            preview = text[:100] + "..." if len(text) > 100 else text
            print(f"  - {preview}")
