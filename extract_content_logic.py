"""Extract content logic and page flow from reference PPT."""
import sys
import io
from pptx import Presentation
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

ref_path = Path(r"C:\Users\18153\Desktop\人工智能应用_马莹奇_总文件夹\01_作品与答辩材料\慧眼识途 - 基于无监督学习的智能交通异常预警系统.pptx")

prs = Presentation(ref_path)

# Content pages mapping
content_pages = {
    1: "封面",
    5: "研究背景 - 政策支持",
    6: "研究背景 - 真实案例",
    7: "研究对象",
    8: "研究现状 - 方案对比",
    9: "研究现状 - 技术瓶颈",
    10: "研究思路 - 核心创新",
    14: "项目流程图",
    15: "模块1 - 边缘数据感知",
    16: "模块2 - 高维表征提取",
    17: "模块3 - 相空间动力学建模",
    18: "模块4 - 多模态语义告警",
    22: "实验结果 - 数据集",
    23: "实验结果 - 实验环境",
    24: "实验结果 - 定量评价",
    25: "实验结果 - 可视化",
    29: "项目总结"
}

print("=" * 100)
print("参考 PPT 内容逻辑与演讲流程")
print("=" * 100)

print("\n【演讲结构】\n")

sections = {
    "第一部分：问题引入": [1, 5, 6, 7],
    "第二部分：现状分析": [8, 9],
    "第三部分：解决方案": [10],
    "第四部分：技术实现": [14, 15, 16, 17, 18],
    "第五部分：实验验证": [22, 23, 24, 25],
    "第六部分：总结": [29]
}

for section, pages in sections.items():
    print(f"{section}")
    for page_num in pages:
        print(f"  第 {page_num} 页: {content_pages[page_num]}")
    print()

print("=" * 100)
print("【每页详细内容提取】")
print("=" * 100)

for page_num, page_title in content_pages.items():
    slide = prs.slides[page_num - 1]
    print(f"\n第 {page_num} 页: {page_title}")
    print("-" * 100)

    # Extract all text
    all_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            # Skip page numbers and section headers
            if text not in ["一. 项目背景", "二. 研究内容", "三. 实验结果", "四.项目总结"] and not text.endswith("/ 15"):
                all_text.append(text)

    # Print main content
    if all_text:
        print("核心内容:")
        for i, text in enumerate(all_text[:5], 1):  # Show first 5 text blocks
            if len(text) > 100:
                print(f"  {i}. {text[:100]}...")
            else:
                print(f"  {i}. {text}")

    # Identify page type
    if page_num == 1:
        print("\n页面类型: 封面")
        print("设计要素: 大标题 + 副标题 + 团队成员 + 背景图")
    elif page_num in [5, 6]:
        print("\n页面类型: 背景介绍")
        print("设计要素: 正文段落 + 配图/案例")
    elif page_num == 7:
        print("\n页面类型: 研究对象")
        print("设计要素: 简短描述 + 示意图")
    elif page_num in [8, 9]:
        print("\n页面类型: 问题分析")
        print("设计要素: 对比图/痛点列举 + 配图")
    elif page_num == 10:
        print("\n页面类型: 解决方案")
        print("设计要素: 三大创新点 + 图标/示意图")
    elif page_num == 14:
        print("\n页面类型: 整体架构")
        print("设计要素: 流程图 + 模块说明")
    elif page_num in [15, 16, 17, 18]:
        print("\n页面类型: 技术细节")
        print("设计要素: 模块名 + 技术说明 + 流程图/代码")
    elif page_num in [22, 23]:
        print("\n页面类型: 实验设置")
        print("设计要素: 数据集/环境说明 + 配图")
    elif page_num in [24, 25]:
        print("\n页面类型: 实验结果")
        print("设计要素: 数据表格/可视化图表")
    elif page_num == 29:
        print("\n页面类型: 总结")
        print("设计要素: 技术链路 + 核心指标")

print("\n" + "=" * 100)
print("【关键设计原则】")
print("=" * 100)

print("""
1. 渐进式叙事:
   - 从宏观到微观（政策 → 案例 → 技术）
   - 从问题到方案（痛点 → 创新 → 实现）
   - 从理论到实践（架构 → 模块 → 实验）

2. 视觉层次:
   - 每页只有 1-2 个核心信息点
   - 文字精简，配图丰富
   - 用颜色/大小区分重要性

3. 技术深度:
   - 背景部分：通俗易懂（政策、案例）
   - 技术部分：专业术语（SigLIP、BLIP、Mahalanobis）
   - 实验部分：数据说话（Cohen's d、检测率）

4. 对比手法:
   - 传统方案 vs 我们的方案（第 8 页）
   - 事后检测 vs 事前预警（第 8 页）
   - 训练集 vs 测试集（第 24 页）

5. 页码编排:
   - 内容页从 1/15 开始
   - 跳过封面、目录、过渡页
   - 总共 15 页内容页

6. 章节划分:
   - 一. 项目背景（6 页）
   - 二. 研究内容（5 页）
   - 三. 实验结果（4 页）
   - 四. 项目总结（1 页）
""")

print("=" * 100)
