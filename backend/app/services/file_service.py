"""Helpers for saving uploads and extracting useful text for AI prompts."""

from __future__ import annotations

from pathlib import Path
from typing import Dict, List
from uuid import uuid4
from xml.etree import ElementTree
from zipfile import BadZipFile, ZipFile

from fastapi import UploadFile
from pptx import Presentation
from pypdf import PdfReader


class FileService:
    """Upload handling for PDF/PPT/PPTX/DOCX files used as resume context."""

    def __init__(self, upload_dir: Path):
        self.upload_dir = upload_dir
        self.upload_dir.mkdir(parents=True, exist_ok=True)

    def _save_file(self, upload: UploadFile) -> Path:
        """Persist an uploaded file to the backend/uploads folder."""

        suffix = Path(upload.filename or "").suffix.lower()
        saved_path = self.upload_dir / f"{uuid4().hex}{suffix}"
        with saved_path.open("wb") as target:
            target.write(upload.file.read())
        upload.file.close()
        return saved_path

    def _extract_pdf_text(self, path: Path) -> str:
        """Extract text from PDF pages."""

        reader = PdfReader(str(path))
        page_text: List[str] = []
        for page in reader.pages:
            page_text.append(page.extract_text() or "")
        return "\n".join(page_text).strip()

    def _extract_pptx_text(self, path: Path) -> str:
        """Extract visible text from a PPTX presentation."""

        presentation = Presentation(str(path))
        slides_text: List[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slides_text.append(shape.text.strip())
        return "\n".join(slides_text).strip()

    def _extract_plain_text(self, path: Path) -> str:
        """Read UTF-8 text files used for direct resume uploads."""

        return path.read_text(encoding="utf-8", errors="ignore").strip()

    def _extract_docx_text(self, path: Path) -> str:
        """Extract visible text from a DOCX document."""

        try:
            with ZipFile(path) as archive:
                xml_bytes = archive.read("word/document.xml")
        except (BadZipFile, KeyError):
            return ""

        root = ElementTree.fromstring(xml_bytes)
        paragraphs: List[str] = []
        namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        for paragraph in root.findall(".//w:p", namespace):
            texts = [node.text or "" for node in paragraph.findall(".//w:t", namespace)]
            merged = "".join(texts).strip()
            if merged:
                paragraphs.append(merged)
        return "\n".join(paragraphs).strip()

    def _extract_text_from_path(self, path: Path) -> tuple[str, str | None]:
        """Extract text from a saved file and return text + optional notice."""

        suffix = path.suffix.lower()
        full_text = ""
        todo_notice = None

        if suffix == ".pdf":
            full_text = self._extract_pdf_text(path)
        elif suffix == ".pptx":
            full_text = self._extract_pptx_text(path)
        elif suffix == ".docx":
            full_text = self._extract_docx_text(path)
        elif suffix in {".txt", ".md"}:
            full_text = self._extract_plain_text(path)
        elif suffix == ".ppt":
            todo_notice = "TODO: legacy .ppt text extraction is not implemented yet."
        elif suffix == ".doc":
            todo_notice = "TODO: legacy .doc text extraction is not implemented yet. Please convert to .docx."
        else:
            todo_notice = "TODO: unsupported file type; only metadata was saved."

        return full_text, todo_notice

    def _resolve_saved_path(self, saved_name: str) -> Path:
        """Resolve one saved upload path inside the upload directory only."""

        safe_name = Path(saved_name or "").name
        if not safe_name or safe_name != (saved_name or ""):
            raise FileNotFoundError("Saved upload was not found.")
        path = self.upload_dir / safe_name
        if not path.exists():
            raise FileNotFoundError("Saved upload was not found.")
        return path

    def preview_saved_file(self, saved_name: str) -> Dict[str, str]:
        """Preview one previously uploaded file by re-extracting its text."""

        path = self._resolve_saved_path(saved_name)
        full_text, todo_notice = self._extract_text_from_path(path)
        return {
            "saved_name": path.name,
            "file_type": path.suffix.lower(),
            "extracted_text_preview": (full_text[:300] + "...") if len(full_text) > 300 else full_text,
            "full_text": full_text,
            "todo_notice": todo_notice,
        }

    def resolve_saved_path(self, saved_name: str) -> Path:
        """Expose a safe saved-upload path lookup for isolated experimental flows."""

        return self._resolve_saved_path(saved_name)

    def delete_saved_file(self, saved_name: str) -> bool:
        """Delete one previously saved upload file from disk."""

        try:
            path = self._resolve_saved_path(saved_name)
        except FileNotFoundError:
            return False

        path.unlink(missing_ok=True)
        return True

    def process_upload(self, upload: UploadFile) -> Dict[str, str]:
        """Save the upload locally and extract a preview when possible.

        Example:
        - Input: PDF named `portfolio.pdf`
        - Output fields include `original_name`, `saved_name`, `extracted_text_preview`, and `full_text`
        """

        saved_path = self._save_file(upload)
        suffix = saved_path.suffix.lower()
        full_text, todo_notice = self._extract_text_from_path(saved_path)

        return {
            "original_name": upload.filename or "unknown_file",
            "saved_name": saved_path.name,
            "file_type": suffix,
            "extracted_text_preview": (full_text[:300] + "...") if len(full_text) > 300 else full_text,
            "full_text": full_text,
            "todo_notice": todo_notice,
        }
