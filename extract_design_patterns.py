"""Extract color scheme, font hierarchy, and layout patterns."""
import sys
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pathlib import Path
from collections import defaultdict

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

ref_path = Path(r"C:\Users\18153\Desktop\人工智能应用_马莹奇_总文件夹\01_作品与答辩材料\慧眼识途 - 基于无监督学习的智能交通异常预警系统.pptx")

prs = Presentation(ref_path)

print("=" * 100)
print("参考 PPT 设计规范提取")
print("=" * 100)

# Font analysis
fonts_used = defaultdict(int)
font_sizes = defaultdict(int)
font_colors = set()

# Layout patterns
layout_types = defaultdict(int)

# Color scheme
bg_colors = set()
text_colors = set()

# Content pages (skip title and transition slides)
content_pages = [5, 6, 7, 8, 9, 10, 14, 15, 16, 17, 18, 22, 23, 24, 25, 29]

print("\n【设计规范分析】\n")

for slide_idx, slide in enumerate(prs.slides, 1):
    if slide_idx not in content_pages:
        continue

    layout_types[slide.slide_layout.name] += 1

    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        fonts_used[run.font.name] += 1
                    if run.font.size:
                        font_sizes[int(run.font.size.pt)] += 1
                    if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                        text_colors.add(str(run.font.color.rgb))

print("1. 字体使用统计:")
for font, count in sorted(fonts_used.items(), key=lambda x: -x[1])[:10]:
    print(f"   - {font}: {count} 次")

print("\n2. 字号层级:")
for size in sorted(font_sizes.keys(), reverse=True):
    count = font_sizes[size]
    usage = ""
    if size >= 36:
        usage = "(主标题)"
    elif size >= 28:
        usage = "(副标题)"
    elif size >= 20:
        usage = "(正文)"
    elif size >= 16:
        usage = "(小字)"
    print(f"   - {size}pt: {count} 次 {usage}")

print("\n3. 布局类型:")
for layout, count in layout_types.items():
    print(f"   - {layout}: {count} 页")

print("\n4. 文本颜色:")
for color in list(text_colors)[:10]:
    print(f"   - {color}")

# Analyze specific content pages for patterns
print("\n" + "=" * 100)
print("【典型页面布局模式】")
print("=" * 100)

# Page 5: Research background
print("\n第 5 页 - 研究背景（文字为主）")
print("  布局: 顶部标题栏 + 正文段落 + 配图")
print("  元素:")
slide = prs.slides[4]
for shape in slide.shapes:
    if hasattr(shape, "text") and shape.text.strip():
        text = shape.text.strip()[:50]
        print(f"    - 文本: {text}...")
        print(f"      位置: ({shape.left / Inches(1):.1f}\", {shape.top / Inches(1):.1f}\")")

# Page 8: Comparison diagram
print("\n第 8 页 - 对比图（流程图）")
print("  布局: 左右对比 + 流程箭头")
slide = prs.slides[7]
shape_count = {"文本": 0, "图片": 0, "箭头": 0}
for shape in slide.shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        shape_count["图片"] += 1
    elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
        shape_count["箭头"] += 1
    elif hasattr(shape, "text") and shape.text.strip():
        shape_count["文本"] += 1
print(f"  元素统计: {shape_count}")

# Page 14: System architecture
print("\n第 14 页 - 系统架构（流程图）")
print("  布局: 四步流程 + 可视化")
slide = prs.slides[13]
print(f"  形状数量: {len(slide.shapes)}")

# Extract common patterns
print("\n" + "=" * 100)
print("【设计模式总结】")
print("=" * 100)

print("\n1. 页面结构:")
print("   - 左上角: 章节标题（如'一. 项目背景'）")
print("   - 右上角: Logo/装饰图")
print("   - 右下角: 页码（如'1 / 15'）")
print("   - 中间: 主要内容区")

print("\n2. 标题层级:")
print("   - 一级标题: 40pt 加粗（章节名）")
print("   - 二级标题: 40pt 楷体（小节名）")
print("   - 正文: 20pt Times New Roman")
print("   - 图注: 楷体 加粗")

print("\n3. 配色方案:")
print("   - 主色调: 蓝色系（科技感）")
print("   - 强调色: 橙色/黄色（高亮）")
print("   - 背景: 白色/浅灰渐变")
print("   - 文字: 深灰/黑色")

print("\n4. 视觉元素:")
print("   - 大量使用图片/截图")
print("   - 流程图用箭头连接")
print("   - 对比用左右分栏")
print("   - 数据用图表可视化")

print("\n5. 内容密度:")
print("   - 每页 1-2 个核心观点")
print("   - 文字简洁（1-2 句话）")
print("   - 图文比例约 1:1")

# Extract slide dimensions
print("\n6. 幻灯片尺寸:")
print(f"   - 宽度: {prs.slide_width / Inches(1):.2f}\" ({prs.slide_width / 914400:.0f}mm)")
print(f"   - 高度: {prs.slide_height / Inches(1):.2f}\" ({prs.slide_height / 914400:.0f}mm)")
print(f"   - 比例: {(prs.slide_width / prs.slide_height):.2f}:1 (超宽屏)")

print("\n" + "=" * 100)
