"""Deep analysis of reference PPT - extract all design patterns."""
import sys
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

ref_path = Path(r"C:\Users\18153\Desktop\人工智能应用_马莹奇_总文件夹\01_作品与答辩材料\慧眼识途 - 基于无监督学习的智能交通异常预警系统.pptx")

prs = Presentation(ref_path)

print("=" * 100)
print("参考 PPT 深度分析报告")
print("=" * 100)

# Analyze presentation-level properties
print("\n【整体属性】")
print(f"总页数: {len(prs.slides)}")
print(f"幻灯片尺寸: {prs.slide_width / Inches(1):.2f}\" x {prs.slide_height / Inches(1):.2f}\"")

# Analyze each slide in detail
for slide_idx, slide in enumerate(prs.slides, 1):
    print("\n" + "=" * 100)
    print(f"【第 {slide_idx} 页详细分析】")
    print("=" * 100)

    # Title
    if slide.shapes.title:
        title = slide.shapes.title
        print(f"\n标题: {title.text}")
        if title.text_frame.paragraphs:
            para = title.text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
                print(f"  - 字体: {run.font.name}")
                print(f"  - 字号: {run.font.size.pt if run.font.size else 'N/A'} pt")
                print(f"  - 颜色: {run.font.color.rgb if hasattr(run.font.color, 'rgb') else 'N/A'}")
                print(f"  - 加粗: {run.font.bold}")
    else:
        print("\n标题: (无)")

    # Layout
    print(f"\n布局类型: {slide.slide_layout.name}")

    # Background
    if slide.background.fill.type:
        print(f"背景类型: {slide.background.fill.type}")

    # Shapes analysis
    print(f"\n形状数量: {len(slide.shapes)}")

    shape_types = {}
    for shape in slide.shapes:
        shape_type = str(shape.shape_type)
        shape_types[shape_type] = shape_types.get(shape_type, 0) + 1

    print("形状类型分布:")
    for stype, count in shape_types.items():
        print(f"  - {stype}: {count}")

    # Text content analysis
    print("\n文本内容:")
    text_count = 0
    for shape_idx, shape in enumerate(slide.shapes):
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            if shape != slide.shapes.title:
                text_count += 1
                print(f"\n  文本框 {text_count}:")
                print(f"    位置: ({shape.left / Inches(1):.2f}\", {shape.top / Inches(1):.2f}\")")
                print(f"    尺寸: {shape.width / Inches(1):.2f}\" x {shape.height / Inches(1):.2f}\"")

                # Show text preview
                preview = text[:150] + "..." if len(text) > 150 else text
                print(f"    内容: {preview}")

                # Font analysis
                if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        print(f"    字体: {run.font.name}")
                        if run.font.size:
                            print(f"    字号: {run.font.size.pt} pt")
                        print(f"    加粗: {run.font.bold}")

    # Images analysis
    print("\n图片/图表:")
    image_count = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            image_count += 1
            print(f"  图片 {image_count}:")
            print(f"    位置: ({shape.left / Inches(1):.2f}\", {shape.top / Inches(1):.2f}\")")
            print(f"    尺寸: {shape.width / Inches(1):.2f}\" x {shape.height / Inches(1):.2f}\"")
        elif shape.shape_type == 3:  # Chart
            image_count += 1
            print(f"  图表 {image_count}:")
            print(f"    位置: ({shape.left / Inches(1):.2f}\", {shape.top / Inches(1):.2f}\")")
            print(f"    尺寸: {shape.width / Inches(1):.2f}\" x {shape.height / Inches(1):.2f}\"")

    if image_count == 0:
        print("  (无图片/图表)")

    # Only analyze first 10 slides in detail to save space
    if slide_idx >= 10:
        print("\n(后续页面省略详细分析，仅显示关键信息)")
        break

print("\n" + "=" * 100)
print("分析完成")
print("=" * 100)
