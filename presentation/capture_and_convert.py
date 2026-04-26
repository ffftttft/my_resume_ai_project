#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
使用 Playwright 截图每一页，然后生成图片版 PPTX
"""
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
import time
import os

def capture_slides():
    """截图每一页幻灯片"""
    screenshots = []

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(
            viewport={'width': 1920, 'height': 1080},
            device_scale_factor=2
        )
        page = context.new_page()

        print('Loading slides...')
        page.goto('http://localhost:3030', wait_until='networkidle')
        time.sleep(2)

        # 获取总页数
        total_slides = page.evaluate('() => __slidev_nav.total')
        print(f'Total slides: {total_slides}')

        # 截图每一页
        for i in range(1, total_slides + 1):
            print(f'Capturing slide {i}/{total_slides}...')

            # 跳转到指定页
            page.goto(f'http://localhost:3030/{i}', wait_until='networkidle')
            time.sleep(1)

            # 截图
            screenshot_path = f'C:/Users/18153/Desktop/my_resume_ai_project/presentation/screenshots/slide_{i:02d}.png'
            os.makedirs(os.path.dirname(screenshot_path), exist_ok=True)
            page.screenshot(path=screenshot_path, full_page=False)
            screenshots.append(screenshot_path)

        browser.close()

    return screenshots

def create_pptx_from_images(screenshots):
    """从截图创建 PPTX"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for screenshot in screenshots:
        # 添加空白幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 插入图片，填满整个幻灯片
        slide.shapes.add_picture(
            screenshot,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

    # 保存
    output_path = 'C:/Users/18153/Desktop/智铸履途.pptx'
    prs.save(output_path)
    print(f'PPTX saved: {output_path}')

def main():
    print('Step 1: Capturing slides as images...')
    screenshots = capture_slides()

    print('Step 2: Creating PPTX from images...')
    create_pptx_from_images(screenshots)

    print('Done!')

if __name__ == '__main__':
    main()
