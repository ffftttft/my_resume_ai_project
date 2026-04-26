#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将 Slidev Markdown 转换为可编辑的 PPTX
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 读取 slides.md
    with open('slides.md', 'r', encoding='utf-8') as f:
        content = f.read()

    # 解析幻灯片
    slides_content = content.split('---\n')[1:]  # 跳过 frontmatter

    for slide_text in slides_content:
        if not slide_text.strip():
            continue

        # 检测布局类型
        if 'layout: cover' in slide_text or 'class: text-center' in slide_text:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
            add_cover_slide(slide, slide_text)
        elif 'layout: two-cols' in slide_text:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_two_cols_slide(slide, slide_text)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_default_slide(slide, slide_text)

    # 保存
    output_path = 'C:/Users/18153/Desktop/智铸履途.pptx'
    prs.save(output_path)
    print(f'PPTX saved to: {output_path}')

def clean_markdown(text):
    """清理 Markdown 标记"""
    # 移除 HTML 标签
    text = re.sub(r'<[^>]+>', '', text)
    # 移除 Markdown 标题符号
    text = re.sub(r'^#+\s+', '', text, flags=re.MULTILINE)
    # 移除加粗
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    # 移除代码块标记
    text = re.sub(r'```[a-z]*\n', '', text)
    text = text.replace('```', '')
    return text.strip()

def add_cover_slide(slide, content):
    """添加封面页"""
    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 12, 41)

    # 提取标题
    title_match = re.search(r'#\s+(.+)', content)
    if title_match:
        title_text = clean_markdown(title_match.group(1))

        # 标题文本框
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title_text

        # 格式化
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(59, 130, 246)

    # 提取副标题
    subtitle_match = re.search(r'##\s+(.+)', content)
    if subtitle_match:
        subtitle_text = clean_markdown(subtitle_match.group(1))

        left = Inches(1)
        top = Inches(4.2)
        width = Inches(8)
        height = Inches(0.8)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = subtitle_box.text_frame
        text_frame.text = subtitle_text

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(241, 245, 249)

def add_two_cols_slide(slide, content):
    """添加两栏布局"""
    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 12, 41)

    # 提取标题
    title_match = re.search(r'#\s+(.+)', content)
    if title_match:
        title_text = clean_markdown(title_match.group(1))

        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title_text

        p = text_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(59, 130, 246)

    # 分割左右内容
    parts = content.split('::right::')
    left_content = parts[0] if len(parts) > 0 else ''
    right_content = parts[1] if len(parts) > 1 else ''

    # 左侧内容
    left_text = clean_markdown(left_content)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(5.5)
    left_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = left_box.text_frame
    text_frame.text = left_text[:500]  # 限制长度
    text_frame.word_wrap = True

    for p in text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(241, 245, 249)

    # 右侧内容
    right_text = clean_markdown(right_content)
    left = Inches(5.2)
    top = Inches(1.5)
    width = Inches(4.3)
    height = Inches(5.5)
    right_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = right_box.text_frame
    text_frame.text = right_text[:500]
    text_frame.word_wrap = True

    for p in text_frame.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(241, 245, 249)

def add_default_slide(slide, content):
    """添加默认布局"""
    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 12, 41)

    # 提取标题
    title_match = re.search(r'#\s+(.+)', content)
    if title_match:
        title_text = clean_markdown(title_match.group(1))

        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.8)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        text_frame.text = title_text

        p = text_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(59, 130, 246)

    # 内容
    body_text = clean_markdown(content)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)
    body_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = body_box.text_frame
    text_frame.text = body_text[:800]
    text_frame.word_wrap = True

    for p in text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(241, 245, 249)

if __name__ == '__main__':
    create_presentation()
