#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
从 PNG 图片创建 PPTX
"""
from pptx import Presentation
from pptx.util import Inches
import os
import glob

def create_pptx_from_images():
    """从截图创建 PPTX"""
    prs = Presentation()
    # 匹配导出图片的 1960x1104 比例
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.633)  # 保持 1960:1104 比例

    # 获取所有 PNG 文件并排序
    screenshots_dir = 'C:/Users/18153/Desktop/my_resume_ai_project/presentation/screenshots'
    screenshots = sorted(glob.glob(f'{screenshots_dir}/*.png'), key=lambda x: int(os.path.basename(x).split('.')[0]))

    print(f'Found {len(screenshots)} slides')

    for i, screenshot in enumerate(screenshots, 1):
        print(f'Adding slide {i}/{len(screenshots)}: {os.path.basename(screenshot)}')

        # 添加空白幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 插入图片，填满整个幻灯片
        slide.shapes.add_picture(
            screenshot,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

    # 保存
    output_path = 'C:/Users/18153/Desktop/智铸履途.pptx'
    prs.save(output_path)
    print(f'PPTX saved: {output_path}')

if __name__ == '__main__':
    create_pptx_from_images()
