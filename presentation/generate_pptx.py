#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
直接生成 PPTX，不依赖转换
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

def set_background_color(slide, rgb):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def add_title_slide(prs):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, RGBColor(15, 12, 41))

    # 主标题
    left, top, width, height = Inches(1), Inches(2.5), Inches(8), Inches(1.2)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "智铸履途"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(59, 130, 246)

    # 副标题
    left, top, width, height = Inches(1), Inches(4), Inches(8), Inches(0.6)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    tf.text = "基于 RAG 与语义 ATS 的智能简历优化系统"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(241, 245, 249)

    # 描述
    left, top, width, height = Inches(1), Inches(4.8), Inches(8), Inches(0.5)
    desc_box = slide.shapes.add_textbox(left, top, width, height)
    tf = desc_box.text_frame
    tf.text = "从原始经历到完美简历，AI 驱动的全栈优化引擎"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(241, 245, 249)

    # 底部信息
    left, top, width, height = Inches(1), Inches(6.5), Inches(8), Inches(0.4)
    info_box = slide.shapes.add_textbox(left, top, width, height)
    tf = info_box.text_frame
    tf.text = "全国大学生计算机程序设计大赛 · 软件应用与开发（Web 应用与开发）"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(241, 245, 249)

    # 汇报人
    left, top, width, height = Inches(1), Inches(7), Inches(8), Inches(0.3)
    author_box = slide.shapes.add_textbox(left, top, width, height)
    tf = author_box.text_frame
    tf.text = "汇报人：陈科源"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(241, 245, 249)

def add_toc_slide(prs):
    """目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, RGBColor(15, 12, 41))

    # 标题
    left, top, width, height = Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "目录"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(59, 130, 246)

    # 四个部分
    sections = [
        ("01", "项目背景", "行业现状 · 核心创新", RGBColor(59, 130, 246)),
        ("02", "核心技术架构", "Web 架构 · RAG 检索", RGBColor(139, 92, 246)),
        ("03", "算法实现与评估", "实验设计 · 性能对比", RGBColor(6, 182, 212)),
        ("04", "项目总结", "技术链路 · 未来展望", RGBColor(16, 185, 129))
    ]

    start_top = 2
    for i, (num, title, subtitle, color) in enumerate(sections):
        row = i // 2
        col = i % 2
        left = Inches(1 + col * 4.5)
        top = Inches(start_top + row * 2.2)
        width, height = Inches(3.5), Inches(1.8)

        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 编号
        p = tf.paragraphs[0]
        p.text = num
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = color

        # 标题
        p = tf.add_paragraph()
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(241, 245, 249)
        p.space_before = Pt(8)

        # 副标题
        p = tf.add_paragraph()
        p.text = subtitle
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(241, 245, 249)
        p.space_before = Pt(4)

def add_section_divider(prs, section_num, title, subtitle):
    """章节过渡页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, RGBColor(15, 12, 41))

    # 标题
    left, top, width, height = Inches(1), Inches(3), Inches(8), Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = f"{section_num}. {title}"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(59, 130, 246)

    # 副标题
    left, top, width, height = Inches(1), Inches(4.5), Inches(8), Inches(0.6)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(241, 245, 249)

def add_content_slide(prs, title, subtitle, content_lines):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, RGBColor(15, 12, 41))

    # 标题
    left, top, width, height = Inches(0.5), Inches(0.4), Inches(9), Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(59, 130, 246)

    # 副标题
    if subtitle:
        left, top, width, height = Inches(0.5), Inches(1.1), Inches(9), Inches(0.4)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        tf = subtitle_box.text_frame
        tf.text = subtitle
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(241, 245, 249)

    # 内容
    left, top, width, height = Inches(0.5), Inches(1.8), Inches(9), Inches(5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(241, 245, 249)
        p.space_after = Pt(8)
        if line.startswith("•") or line.startswith("-"):
            p.level = 1

def add_final_slide(prs):
    """结尾页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, RGBColor(15, 12, 41))

    # 主标题
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "智铸履途"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(59, 130, 246)

    # 副标题
    left, top, width, height = Inches(1), Inches(3.5), Inches(8), Inches(0.6)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    tf.text = "锻造完美简历，铸就职场未来"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(241, 245, 249)

    # Q&A
    left, top, width, height = Inches(1), Inches(4.5), Inches(8), Inches(0.5)
    qa_box = slide.shapes.add_textbox(left, top, width, height)
    tf = qa_box.text_frame
    tf.text = "Q & A"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(241, 245, 249)

    # 核心指标
    metrics = [
        ("87%", "ATS 匹配度", RGBColor(16, 185, 129)),
        ("<500ms", "首字节响应", RGBColor(59, 130, 246)),
        ("100%", "结构化契约", RGBColor(139, 92, 246))
    ]

    for i, (value, label, color) in enumerate(metrics):
        left = Inches(1.5 + i * 2.5)
        top = Inches(5.5)
        width, height = Inches(2), Inches(1)

        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame

        p = tf.paragraphs[0]
        p.text = value
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color

        p = tf.add_paragraph()
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(241, 245, 249)
        p.space_before = Pt(4)

    # 底部
    left, top, width, height = Inches(1), Inches(6.8), Inches(8), Inches(0.4)
    footer_box = slide.shapes.add_textbox(left, top, width, height)
    tf = footer_box.text_frame
    tf.text = "欢迎各位评委批评指正"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(241, 245, 249)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 封面
    add_title_slide(prs)

    # 2. 目录
    add_toc_slide(prs)

    # 3. 第一部分：项目背景
    add_section_divider(prs, "一", "项目背景", "行业现状 · 核心创新 · 研究目标")

    add_content_slide(prs, "背景挑战", "传统 AI 简历生成的黑盒困局", [
        "通用 LLM 缺乏行业领域知识，且生成的 JSON 格式极不稳定。",
        "本项目通过 RAG 范式引入专家知识库，解决 Web 端内容生成的不可控难题。",
        "",
        "传统工具痛点：",
        "• 缺乏领域知识，内容空洞",
        "• JSON 格式不稳定，解析失败",
        "• 黑盒生成，缺乏可解释性",
        "",
        "行业数据：",
        "• 75% 简历被 ATS 系统过滤",
        "• HR 平均阅读时间仅 3 秒",
        "• 2024 届毕业生人数达 1179 万"
    ])

    add_content_slide(prs, "核心创新", "构建高可解释性的结构化简历生成闭环", [
        "我们不仅在做 Web 应用，更在定义一套基于 Pydantic 契约的简历生成标准，",
        "确保生成质量的 100% 完整性与可编辑性。",
        "",
        "四大核心技术：",
        "• RAG 检索：零标注领域知识注入",
        "• 语义 ATS：1536 维向量空间匹配",
        "• 结构化输出：Pydantic 契约保证",
        "• 流式传输：SSE 实时反馈",
        "",
        "核心目标：从通用模板到个性化定制",
        "AI 理解用户背景 → 检索优质案例 → 生成针对性简历 → 语义评分优化"
    ])

    add_content_slide(prs, "多模态输入", "PDF 与 PPTX 全自动解析", [
        "系统支持多模态输入，能够从项目 PPT 中捕捉用户被隐藏的硬技能特征。",
        "",
        "支持格式：",
        "• PDF 简历：pypdf 提取",
        "• PPTX 项目：python-pptx 解析",
        "• 文本分片：语义切分",
        "• 自动向量化：Embedding API",
        "",
        "技术优势：",
        "• 自动提取项目经验",
        "• 识别隐藏技能点",
        "• 丰富原始特征"
    ])

    add_content_slide(prs, "ProfileMemory", "4KB 压缩记忆与主动追问", [
        "仿照人类专家的对话逻辑，当信息不足时主动发起追问。",
        "",
        "核心功能：",
        "• 自动识别信息缺失",
        "• 生成针对性追问",
        "• 压缩历史对话（4KB）",
        "• 多轮上下文保持",
        "",
        "记忆压缩策略：",
        "保留最近 N 条对话 + 关键信息摘要"
    ])

    add_content_slide(prs, "RAG 检索增强", "ChromaDB 驱动的专家知识库", [
        "通过 RAG 技术，将行业专家的撰写规范动态注入 Prompt。",
        "",
        "RAG 优势：",
        "• 零标注数据：快速部署，无需训练",
        "• 动态知识更新：无需重训练模型",
        "• 领域专家规范：注入行业最佳实践",
        "• 可解释结果：检索过程透明",
        "",
        "数据来源：",
        "• 中文参考简历库：50+ 条",
        "• 英文参考简历库：50+ 条",
        "• 持续更新中"
    ])

    add_content_slide(prs, "结构化输出", "Instructor + Pydantic 契约保证", [
        "摒弃传统 json.loads 解析，采用 Pydantic v2.0 强制验证。",
        "",
        "传统 JSON 解析的问题：",
        "• 不稳定，容易出错",
        "• 可能缺少字段",
        "• 可能类型错误",
        "• 解析失败需要重试",
        "",
        "结构化输出的优势：",
        "• 100% 结构完整性",
        "• 类型安全保证",
        "• 自动验证和修复",
        "• 完美支持 Web 应用"
    ])

    add_content_slide(prs, "SSE 流式传输", "Server-Sent Events 实时反馈", [
        "针对 Web 应用场景，实现 SSE 异步传输，首字节时间 < 500ms。",
        "",
        "性能对比：",
        "• 传统模式：8 秒",
        "• 流式模式：3 秒",
        "• 体验提升：62%",
        "",
        "技术实现：",
        "• 后端：FastAPI StreamingResponse",
        "• 前端：EventSource 接收",
        "• 实时更新预览界面",
        "• 首字节响应 < 500ms"
    ])

    add_content_slide(prs, "语义 ATS 评分", "Embedding 相似度量化匹配", [
        "基于 1536 维特征的语义距离，输出详细多维度评估报告。",
        "",
        "评分维度：",
        "• 关键词覆盖率：JD 关键词匹配度",
        "• 技能匹配度：技术栈相似度",
        "• 经验相关性：工作经历匹配",
        "• 格式规范性：ATS 友好度",
        "• 语义相似度：Cosine Similarity",
        "",
        "综合评分算法：",
        "overall = (coverage * 0.4 + semantic_score * 0.6) * 100"
    ])

    # 4. 第二部分：核心技术架构
    add_section_divider(prs, "二", "核心技术架构", "Web 架构 · RAG 检索 · 结构化生成")

    add_content_slide(prs, "系统全景", "基于微服务解耦的高性能 Web 架构", [
        "我们采用了异步 FastAPI 框架与 Vite 构建工具，",
        "确保了首字节时间小于 500ms 的极致性能表现。",
        "",
        "架构分层：",
        "• 前端层：React 19 + Vite 构建 + SSE 流式接收",
        "• API 层：FastAPI 路由 + 服务编排",
        "• 服务层：ProfileMemory + RAG + Resume + ATS",
        "• 数据层：ChromaDB 向量数据库 + JSON 参考库",
        "",
        "技术栈：",
        "FastAPI · React 19 · ChromaDB · Pydantic · SSE"
    ])

    # 5. 第三部分：算法实现与评估
    add_section_divider(prs, "三", "算法实现与评估", "实验设计 · 性能对比 · 可视化分析")

    add_content_slide(prs, "性能领先", "核心质量指标与生成效率的全面胜出", [
        "实验证明，本系统在各方面均显著优于通用大模型，",
        "尤其是语义匹配度和结构完整性指标达到了商用级水平。",
        "",
        "对比数据（智简 AI vs 豆包 vs ChatGPT）：",
        "• ATS 匹配度：87% vs 65% vs 72%（+22%）",
        "• 关键词覆盖：18/20 vs 12/20 vs 15/20（+50%）",
        "• 结构完整性：5/5 vs 3/5 vs 4/5（完美）",
        "• 生成时间：3s vs 8s vs 6s（-62%）",
        "• 首字节时间：<500ms（流式优势）",
        "",
        "核心优势：RAG 领域知识注入 + Pydantic 契约保证 + SSE 流式传输"
    ])

    # 6. 第四部分：项目总结
    add_section_divider(prs, "四", "项目总结", "技术链路 · 核心指标 · 未来展望")

    add_content_slide(prs, "全栈闭环", "打造具备 Git 级版本控制的简历生态", [
        "我们成功构建了从感知到量化评分的完整 Web 技术闭环。",
        "",
        "技术链路：",
        "• 追问式对话 → ProfileMemory（4KB）",
        "• RAG 检索 → ChromaDB 知识注入",
        "• 结构化生成 → Pydantic 契约",
        "• 语义 ATS → 1536 维 Embedding",
        "• 流式输出 → SSE 实时反馈",
        "",
        "核心指标：",
        "• ATS 匹配度：87%（+22%）",
        "• 关键词覆盖：18/20（+50%）",
        "• 结构完整性：5/5（完美）",
        "• 生成速度：3s（-62%）",
        "• 首字节：<500ms",
        "• 支持中英文双语",
        "",
        "未来展望：多轮对话优化 · 职位推荐 · Git-like 版本管理 · 面试准备辅助"
    ])

    # 7. 结尾页
    add_final_slide(prs)

    # 保存
    output_path = 'C:/Users/18153/Desktop/智铸履途.pptx'
    prs.save(output_path)
    print(f'PPTX saved: {output_path}')

if __name__ == '__main__':
    create_presentation()
